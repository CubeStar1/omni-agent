import os
import tempfile
import io
from typing import List, Optional, Literal
from langchain_core.documents import Document
from langchain_markitdown import PptxLoader as BasePptxLoader
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import cv2
import numpy as np
import easyocr
class CustomPptxLoader:
    """
    Custom PPTX loader that can extract text using standard text extraction
    or OCR from slide images.
    """
    
    def __init__(self, file_path: str, use_ocr: bool = False, ocr_engine: Literal["pytesseract", "easyocr"] = "easyocr"):
        """
        Initialize the custom PPTX loader.
        
        Args:
            file_path (str): Path to the PPTX file
            use_ocr (bool): Whether to use OCR for text extraction
            ocr_engine (str): OCR engine to use ("pytesseract" or "easyocr")
        """
        self.file_path = file_path
        self.use_ocr = use_ocr
        self.ocr_engine = ocr_engine
        
        self._easyocr_reader = None
        if self.use_ocr and self.ocr_engine == "easyocr":
            print("Initializing EasyOCR reader...")
            self._easyocr_reader = easyocr.Reader(['en'], gpu=True) 
            print("EasyOCR reader initialized")
    
    def load(self) -> List[Document]:
        """
        Load the PPTX file and extract text.
        
        Returns:
            List[Document]: List of documents with extracted text
        """
        if self.use_ocr:
            return self._load_with_ocr()
        else:
            return self._load_standard()
    
    def _load_standard(self) -> List[Document]:
        """
        Load PPTX using standard text extraction (langchain's PptxLoader).
        
        Returns:
            List[Document]: List of documents with extracted text
        """
        try:
            loader = BasePptxLoader(self.file_path)
            documents = loader.load()
            
            for doc in documents:
                doc.metadata.update({
                    'extraction_method': 'standard',
                    'file_type': 'pptx'
                })
            
            return documents
        except Exception as e:
            raise Exception(f"Failed to load PPTX with standard extraction: {str(e)}")
    
    def _load_with_ocr(self) -> List[Document]:
        """
        Load PPTX using OCR by extracting images directly from the presentation.
        
        Returns:
            List[Document]: List of documents with OCR-extracted text
        """
        documents = []
        
        try:
            print("Loading PPTX presentation and extracting images...")
            prs = Presentation(self.file_path)
            
            images_with_metadata = self._extract_images_from_pptx(prs)
            
            if not images_with_metadata:
                print("⚠️ No images found in presentation, falling back to standard text extraction")
                return self._load_standard()
            
            print(f"🔍 Found {len(images_with_metadata)} images, extracting text using {self.ocr_engine.upper()}...")
            
            for img_data in images_with_metadata:
                try:
                    image = img_data['image']
                    slide_num = img_data['slide_number']
                    image_index = img_data['image_index']
                    
                    text = self._extract_text_from_image(image)
                    
                    if text.strip():
                        doc = Document(
                            page_content=text.strip(),
                            metadata={
                                'source': self.file_path,
                                'slide_number': slide_num,
                                'image_index': image_index,
                                'extraction_method': f'OCR_{self.ocr_engine}',
                                'file_type': 'pptx'
                            }
                        )
                        documents.append(doc)
                    else:
                        doc = Document(
                            page_content=f"[Slide {slide_num}, Image {image_index} - No readable text detected]",
                            metadata={
                                'source': self.file_path,
                                'slide_number': slide_num,
                                'image_index': image_index,
                                'extraction_method': f'OCR_{self.ocr_engine}',
                                'file_type': 'pptx',
                                'no_text': True
                            }
                        )
                        documents.append(doc)
                        
                except Exception as e:
                    print(f"⚠️ Failed to extract text from slide {slide_num}, image {image_index}: {str(e)}")
                    doc = Document(
                        page_content=f"[Slide {slide_num}, Image {image_index} - Error extracting text: {str(e)}]",
                        metadata={
                            'source': self.file_path,
                            'slide_number': slide_num,
                            'image_index': image_index,
                            'extraction_method': f'OCR_{self.ocr_engine}',
                            'file_type': 'pptx',
                            'error': str(e)
                        }
                    )
                    documents.append(doc)
            
            if not documents:
                print("No text could be extracted from images, falling back to standard text extraction")
                return self._load_standard()
            
            try:
                standard_docs = self._load_standard()
                for doc in standard_docs:
                    doc.metadata['extraction_method'] = 'standard_combined'
                documents.extend(standard_docs)
            except Exception as e:
                print(f"Could not extract standard text: {str(e)}")
            
            print(f"Successfully extracted text from {len([d for d in documents if not d.metadata.get('no_text') and not d.metadata.get('error')])} sources")
            
            return documents
            
        except Exception as e:
            raise Exception(f"Failed to load PPTX with OCR: {str(e)}")
    
    def _extract_text_from_image(self, image: Image.Image) -> str:
        """
        Extract text from an image using the specified OCR engine.
        
        Args:
            image (Image.Image): PIL Image object
            
        Returns:
            str: Extracted text
        """
        try:
            if self.ocr_engine == "pytesseract":
                return pytesseract.image_to_string(image, config='--psm 6')
            elif self.ocr_engine == "easyocr":
                if self._easyocr_reader is None:
                    raise RuntimeError("EasyOCR reader not initialized")
                
                image_np = np.array(image)
                
                if len(image_np.shape) == 3 and image_np.shape[2] == 3:
                    pass
                elif len(image_np.shape) == 3 and image_np.shape[2] == 4:
                    image_np = image_np[:, :, :3]
                
                results = self._easyocr_reader.readtext(image_np, detail=0)
                
                return ' '.join(results)
            else:
                raise ValueError(f"Unsupported OCR engine: {self.ocr_engine}")
                
        except Exception as e:
            raise Exception(f"Failed to extract text using {self.ocr_engine}: {str(e)}")
    
    def _extract_images_from_pptx(self, prs: Presentation) -> List[dict]:
        """
        Extract all images from the PowerPoint presentation.
        
        Args:
            prs (Presentation): The PowerPoint presentation object
            
        Returns:
            List[dict]: List of dictionaries containing image data and metadata
        """
        images_data = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            image_idx = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        nparr = np.frombuffer(image_bytes, np.uint8)
                        cv_image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                        
                        rgb_image = cv2.cvtColor(cv_image, cv2.COLOR_BGR2RGB)
                        
                        pil_image = Image.fromarray(rgb_image)
                        
                        images_data.append({
                            'image': pil_image,
                            'slide_number': slide_idx,
                            'image_index': image_idx,
                            'shape': shape
                        })
                        
                        image_idx += 1
                        
                    except Exception as e:
                        print(f"Failed to extract image from slide {slide_idx}: {str(e)}")
                        continue
        
        return images_data

def load_pptx_with_options(file_path: str, use_ocr: bool = False, ocr_engine: Literal["pytesseract", "easyocr"] = "pytesseract") -> List[Document]:
    """
    Convenience function to load PPTX with options.
    
    Args:
        file_path (str): Path to the PPTX file
        use_ocr (bool): Whether to use OCR for text extraction
        ocr_engine (str): OCR engine to use ("pytesseract" or "easyocr")
        
    Returns:
        List[Document]: List of documents with extracted text
    """
    loader = CustomPptxLoader(file_path, use_ocr=use_ocr, ocr_engine=ocr_engine)
    return loader.load()
